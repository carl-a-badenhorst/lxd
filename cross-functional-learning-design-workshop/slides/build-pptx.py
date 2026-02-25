"""
Generates workshop-slides.pptx for the cross-functional learning design workshop.

Mirrors the 37-slide PDF deck: same content and structure, minimalist utilitarian design.
Slide size: 16:9. Run from slides folder: python build-pptx.py

Image placeholders: right-click shape -> Change Picture. Assets in ../assets/images/
  cover-workshop.jpg, phase-1-framing.jpg, phase-2-science.jpg, phase-3-roles.jpg,
  phase-4-diagnosis.jpg, phase-5-redesign.jpg, close-workshop.jpg
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Design tokens (match PDF)
INK = RGBColor(17, 17, 17)
MID = RGBColor(68, 68, 68)
LIGHT = RGBColor(136, 136, 136)
BG = RGBColor(245, 245, 245)
HIGHLIGHT = RGBColor(232, 232, 232)
WHITE = RGBColor(255, 255, 255)
DARK_BG = RGBColor(17, 17, 17)

M = 0.5   # margin inches
TITLE_PT = 24
BODY_PT = 12
SMALL_PT = 10
LABEL_PT = 9

TOTAL_SLIDES = 37


def _textbox(slide, left, top, width, height, text, font_size=BODY_PT, bold=False, italic=False, color=INK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = "Calibri"
    return box


def _add_para(tf, text, font_size=BODY_PT, bold=False, italic=False, color=INK):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = "Calibri"


def _placeholder_rect(slide, left, top, width, height, label="[ IMAGE ]"):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(222, 222, 222)
    shape.line.color.rgb = LIGHT
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.paragraphs[0].text = label
    tf.paragraphs[0].font.size = Pt(8)
    tf.paragraphs[0].font.color.rgb = LIGHT
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def _footer(slide, n):
    slide.shapes.add_textbox(Inches(M), Inches(7.0), Inches(10), Inches(0.3)).text_frame.paragraphs[0].text = (
        "Designing Learning as a Cross-Functional Activity System"
    )
    slide.shapes.add_textbox(Inches(12), Inches(7.0), Inches(1), Inches(0.3)).text_frame.paragraphs[0].text = f"{n} / {TOTAL_SLIDES}"


def _phase_header_slide(prs, number, title, description, n, image_placeholder=False):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    # Background
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), Inches(7.5))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = DARK_BG
    left_bar.line.fill.background()
    # Grey number
    box = slide.shapes.add_textbox(Inches(M + 0.2), Inches(2.8), Inches(2), Inches(1.2))
    box.text_frame.paragraphs[0].text = number
    box.text_frame.paragraphs[0].font.size = Pt(96)
    box.text_frame.paragraphs[0].font.bold = True
    box.text_frame.paragraphs[0].font.color.rgb = RGBColor(224, 224, 224)
    # PHASE label
    _textbox(slide, M + 0.2, 0.4, 2, 0.25, "PHASE", LABEL_PT, bold=True, color=LIGHT)
    # Title
    _textbox(slide, M + 0.2, 0.7, 7, 0.6, title, 28, bold=True, color=INK)
    # Description
    desc_box = slide.shapes.add_textbox(Inches(M + 0.2), Inches(1.5), Inches(7), Inches(1.2))
    desc_box.text_frame.word_wrap = True
    desc_box.text_frame.paragraphs[0].text = description
    desc_box.text_frame.paragraphs[0].font.size = Pt(BODY_PT)
    desc_box.text_frame.paragraphs[0].font.color.rgb = MID
    if image_placeholder:
        _placeholder_rect(slide, 8.2, 1.2, 4.5, 5, "[ IMAGE ]")
    _footer(slide, n)


def build_pptx(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # --- Slide 1: Title ---
    s = prs.slides.add_slide(blank)
    left_dark = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(8), Inches(7.5))
    left_dark.fill.solid()
    left_dark.fill.fore_color.rgb = DARK_BG
    left_dark.line.fill.background()
    white_panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), 0, Inches(5.333), Inches(7.5))
    white_panel.fill.solid()
    white_panel.fill.fore_color.rgb = WHITE
    white_panel.line.fill.background()
    _placeholder_rect(s, 0, 0, 8, 7.5, "[ COVER IMAGE ]")
    _textbox(s, 8.2, 0.4, 4.8, 0.3, "WORKSHOP", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, 8.2, 0.8, 4.8, 0.5, "Designing Learning", 22, bold=True, color=INK)
    _textbox(s, 8.2, 1.25, 4.8, 0.5, "as a Cross-Functional", 22, bold=True, color=INK)
    _textbox(s, 8.2, 1.7, 4.8, 0.5, "Activity System", 22, bold=True, color=INK)
    _textbox(s, 8.2, 2.5, 4.8, 0.3, "135-minute design workshop", SMALL_PT, color=MID)
    s.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.6), Inches(0.3)).text_frame.paragraphs[0].text = f"1 / {TOTAL_SLIDES}"

    # --- Slide 2: Agenda ---
    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 2, 0.25, "OVERVIEW", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 8, 0.5, "Workshop agenda", TITLE_PT, bold=True, color=INK)
    agenda_data = [
        ("Phase 1", "0\u201315 min", "Framing the Shared Object"),
        ("Phase 2", "15\u201330 min", "Learning Science Core"),
        ("Phase 3", "30\u201350 min", "Role Mapping"),
        ("Phase 4", "50\u201375 min", "Cross-Role Diagnosis"),
        ("\u2014", "75\u201385 min", "Break"),
        ("Phase 5", "85\u2013110 min", "Collaborative Redesign"),
        ("Phase 6", "110\u2013125 min", "Cross-Team Critique"),
        ("Phase 7", "125\u2013135 min", "Collaboration Charter & Closing"),
    ]
    tbl = s.shapes.add_table(9, 3, Inches(M), Inches(1.15), Inches(12), Inches(2.8)).table
    tbl.cell(0, 0).text = "Phase"
    tbl.cell(0, 1).text = "Time"
    tbl.cell(0, 2).text = "Title"
    for i, (phase, time_val, title) in enumerate(agenda_data, 1):
        tbl.cell(i, 0).text = phase
        tbl.cell(i, 1).text = time_val
        tbl.cell(i, 2).text = title
    _footer(s, 2)

    # --- Slide 3: Diagnostic norm ---
    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 2, 0.25, "GROUND RULE", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 8, 0.5, "Before we begin", TITLE_PT, bold=True, color=INK)
    _textbox(s, M, 1.4, 11, 0.4, "\u201cWe\u2019re diagnosing the system, not individuals.", 15, italic=True, color=INK)
    _textbox(s, M, 1.85, 11, 0.4, "Misalignment is structural, not personal.", 15, italic=True, color=INK)
    _textbox(s, M, 2.3, 11, 0.4, "Our goal is coherence, not blame.\u201d", 15, italic=True, color=INK)
    _textbox(s, M, 3.0, 11, 0.35, "Return to this norm whenever tension rises during diagnosis.", SMALL_PT, color=LIGHT)
    _footer(s, 3)

    # --- Slides 4\u20139: Phase 1 ---
    _phase_header_slide(prs, "1", "Framing the Shared Object",
        "We begin by surfacing what we already believe learning is \u2014 then replacing those beliefs with an operational definition.", 4, image_placeholder=True)

    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 3, 0.25, "Phase 1 \u2014 Activity", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 8, 0.5, "Writing prompt 1A", TITLE_PT, bold=True, color=INK)
    callout = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(M), Inches(2.6), Inches(10), Inches(0.6))
    callout.fill.solid()
    callout.fill.fore_color.rgb = HIGHLIGHT
    callout.line.color.rgb = LIGHT
    _textbox(s, M + 0.15, 2.65, 9.5, 0.55, "Define learning in one sentence.", BODY_PT, italic=True, color=INK)
    _textbox(s, M, 3.5, 11, 0.3, "Write without looking at the glossary. 90 seconds. No discussion.", SMALL_PT, color=MID)
    _footer(s, 5)

    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 3, 0.25, "Phase 1 \u2014 Activity", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 8, 0.5, "Writing prompt 1B", TITLE_PT, bold=True, color=INK)
    callout2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(M), Inches(2.6), Inches(10), Inches(0.6))
    callout2.fill.solid()
    callout2.fill.fore_color.rgb = HIGHLIGHT
    _textbox(s, M + 0.15, 2.65, 9.5, 0.55, "How do you know when learning has happened?", BODY_PT, italic=True, color=INK)
    _textbox(s, M, 3.5, 11, 0.3, "One sentence. Write, don\u2019t type. 90 seconds.", SMALL_PT, color=MID)
    _footer(s, 6)

    s = prs.slides.add_slide(blank)
    left_dark = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(6.2), Inches(7.5))
    left_dark.fill.solid()
    left_dark.fill.fore_color.rgb = DARK_BG
    left_dark.line.fill.background()
    _textbox(s, M, 0.35, 5.5, 0.25, "THE OPERATIONAL DEFINITION", LABEL_PT, bold=True, color=RGBColor(136, 136, 136))
    _textbox(s, M, 1.8, 5.5, 0.5, "\u201cLearning is durable change", 18, italic=True, color=WHITE)
    _textbox(s, M, 2.3, 5.5, 0.5, "in knowledge structures", 18, italic=True, color=WHITE)
    _textbox(s, M, 2.8, 5.5, 0.5, "that enables future participation", 18, italic=True, color=WHITE)
    _textbox(s, M, 3.3, 5.5, 0.5, "and performance.\u201d", 18, italic=True, color=WHITE)
    _textbox(s, 6.5, 0.35, 6, 0.25, "THREE WORDS DOING THE WORK", SMALL_PT, bold=True, color=MID)
    _textbox(s, 6.5, 0.75, 6, 0.35, "Durable", BODY_PT, bold=True, color=INK)
    _textbox(s, 6.5, 1.05, 6, 0.5, "Not fleeting. Visible months later, not minutes after.", SMALL_PT, color=MID)
    _textbox(s, 6.5, 1.7, 6, 0.35, "Knowledge structures", BODY_PT, bold=True, color=INK)
    _textbox(s, 6.5, 2.0, 6, 0.5, "Organised schemas in long-term memory, not isolated facts.", SMALL_PT, color=MID)
    _textbox(s, 6.5, 2.65, 6, 0.35, "Enables future", BODY_PT, bold=True, color=INK)
    _textbox(s, 6.5, 2.95, 6, 0.5, "The test is what learners can do later, not during the session.", SMALL_PT, color=MID)
    s.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.6), Inches(0.3)).text_frame.paragraphs[0].text = f"7 / {TOTAL_SLIDES}"

    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 3, 0.25, "Phase 1 \u2014 Insight", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 8, 0.5, "Participation \u2260 learning", TITLE_PT, bold=True, color=INK)
    _textbox(s, M, 0.95, 11, 0.4, "Most definitions of learning describe activity. Activity is what we hope causes learning.", BODY_PT, color=MID)
    _textbox(s, M, 1.5, 3.5, 0.25, "Activity completion", SMALL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 1.85, 3.5, 0.25, "\u2715  Watched the video", SMALL_PT, color=LIGHT)
    _textbox(s, M, 2.1, 3.5, 0.25, "\u2715  Completed the quiz", SMALL_PT, color=LIGHT)
    _textbox(s, M, 2.35, 3.5, 0.25, "\u2715  Attended the session", SMALL_PT, color=LIGHT)
    _textbox(s, M, 2.6, 3.5, 0.25, "\u2715  Submitted the reflection", SMALL_PT, color=LIGHT)
    _textbox(s, 6.8, 1.5, 3.5, 0.25, "Durable capability", SMALL_PT, bold=True, color=INK)
    _textbox(s, 6.8, 1.85, 3.5, 0.25, "\u2713  Recalls and applies the concept 6 months later", SMALL_PT, color=INK)
    _textbox(s, 6.8, 2.15, 3.5, 0.25, "\u2713  Transfers knowledge to a novel situation", SMALL_PT, color=INK)
    _textbox(s, 6.8, 2.45, 3.5, 0.25, "\u2713  Performs competently under real conditions", SMALL_PT, color=INK)
    _textbox(s, 6.8, 2.75, 3.5, 0.25, "\u2713  Judgment improves over time", SMALL_PT, color=INK)
    _footer(s, 8)

    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 3, 0.25, "Phase 1 \u2014 Activity", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 8, 0.5, "Writing prompt 1C", TITLE_PT, bold=True, color=INK)
    _textbox(s, M, 1.4, 11, 0.6, "1.  What must learners be able to do 6\u201312 months after your module, in real conditions?", BODY_PT, color=INK)
    _textbox(s, M, 2.2, 11, 0.6, "2.  What cognitive change must occur for that to be possible?", BODY_PT, color=INK)
    _textbox(s, M, 3.2, 11, 0.3, "3 minutes individual writing. Share two or three responses before moving on.", SMALL_PT, color=MID)
    _footer(s, 9)

    # --- Phase 2 ---
    _phase_header_slide(prs, "2", "Learning Science Core",
        "Three commitments that every design decision in this room must be accountable to.", 10, image_placeholder=True)

    s = prs.slides.add_slide(blank)
    dark_bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    dark_bg.fill.solid()
    dark_bg.fill.fore_color.rgb = DARK_BG
    dark_bg.line.fill.background()
    _textbox(s, M, 0.35, 3, 0.25, "PHASE 2 \u2014 RECALL TASK", LABEL_PT, bold=True, color=RGBColor(136, 136, 136))
    _textbox(s, M, 0.7, 8, 0.5, "Close your notes.", 28, bold=True, color=WHITE)
    _textbox(s, M, 1.8, 11, 0.5, "Without looking, write the three learning science commitments from the pre-work. Two minutes.", 16, italic=True, color=RGBColor(204, 204, 204))
    _textbox(s, M, 3.2, 11, 0.3, "Silence is intentional. Do not help your neighbour.", SMALL_PT, color=RGBColor(102, 102, 102))
    nb = s.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.6), Inches(0.3))
    nb.text_frame.paragraphs[0].text = "11 / 37"
    nb.text_frame.paragraphs[0].font.color.rgb = RGBColor(136, 136, 136)

    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 3, 0.25, "Phase 2 \u2014 Reveal", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 8, 0.5, "The three commitments", TITLE_PT, bold=True, color=INK)
    _textbox(s, M, 1.3, 11, 0.45, "1.  Memory precedes complex thinking", 20, bold=True, color=INK)
    _textbox(s, M, 1.95, 11, 0.45, "2.  Retrieval strengthens memory", 20, bold=True, color=INK)
    _textbox(s, M, 2.6, 11, 0.45, "3.  Cognitive load must be managed", 20, bold=True, color=INK)
    _textbox(s, M, 3.5, 11, 0.4, "How many did you get? These three commitments are the arbitration language for every design decision today.", SMALL_PT, color=LIGHT)
    _footer(s, 12)

    def _commitment_slide(prs, num, title, body, implication, slide_n):
        s = prs.slides.add_slide(blank)
        _textbox(s, M, 0.35, 3, 0.25, f"Commitment {num}", LABEL_PT, bold=True, color=LIGHT)
        _textbox(s, M, 0.55, 9, 0.5, title, TITLE_PT, bold=True, color=INK)
        box = s.shapes.add_textbox(Inches(M), Inches(1.0), Inches(11), Inches(1.8))
        box.text_frame.word_wrap = True
        box.text_frame.paragraphs[0].text = body
        box.text_frame.paragraphs[0].font.size = Pt(BODY_PT)
        box.text_frame.paragraphs[0].font.color.rgb = INK
        imp_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(M), Inches(3.0), Inches(12), Inches(1.2))
        imp_box.fill.solid()
        imp_box.fill.fore_color.rgb = HIGHLIGHT
        imp_box.line.fill.background()
        _textbox(s, M + 0.15, 3.05, 11.5, 1.1, "Design implication: " + implication, SMALL_PT, italic=True, color=MID)
        _footer(s, slide_n)

    _commitment_slide(prs, "1", "Memory precedes complex thinking",
        "You cannot analyse, evaluate or create with knowledge you cannot recall. Higher-order cognitive operations require accessible knowledge structures in long-term memory.",
        "Sequence instruction so foundational knowledge is encoded and retrievable before learners are asked to apply it in complex ways.", 13)
    _commitment_slide(prs, "2", "Retrieval strengthens memory",
        "Actively recalling information from memory strengthens the neural pathways that make future recall easier. Re-reading or re-watching does not produce the same effect.",
        "Build retrieval opportunities into the design with specific timing, mechanisms and spacing. Retrieval is not a test event \u2014 it is an instructional strategy.", 14)
    _commitment_slide(prs, "3", "Cognitive load must be managed",
        "Working memory is limited. Extraneous cognitive load \u2014 effort not contributing to schema formation \u2014 reduces the capacity available for learning. Complexity and novelty must be introduced deliberately, not simultaneously.",
        "Remove presentation elements that do not serve the learning objective. Manage the number of new concepts introduced at once. Use worked examples before independent problem-solving.", 15)

    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 4, 0.25, "Phase 2 \u2014 Role-based application", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 9, 0.5, "Apply the commitments to your work", TITLE_PT, bold=True, color=INK)
    _textbox(s, M, 0.95, 11, 0.3, "Form role-based groups. Discuss the three questions below. 10 minutes.", SMALL_PT, color=MID)
    _textbox(s, M, 1.45, 11, 0.4, "1.  Where is retrieval structured in your typical design work? Name a specific artefact or moment.", SMALL_PT, color=INK)
    _textbox(s, M, 2.0, 11, 0.4, "2.  Where might cognitive load be excessive? Name a specific module, assessment or platform.", SMALL_PT, color=INK)
    _textbox(s, M, 2.55, 11, 0.4, "3.  Where is thinking hidden? Where do learners perform activity without their reasoning being visible?", SMALL_PT, color=INK)
    _footer(s, 16)

    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 3, 0.25, "Phase 2 \u2014 Worked example", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 8, 0.5, "Diagnosis in practice", TITLE_PT, bold=True, color=INK)
    _textbox(s, M, 0.95, 11, 0.35, "Before we diagnose your work, let\u2019s practice on this case.", SMALL_PT, color=MID)
    box1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(M), Inches(2.2), Inches(5.8), Inches(0.9))
    box1.fill.solid()
    box1.fill.fore_color.rgb = BG
    box1.line.color.rgb = LIGHT
    _textbox(s, M + 0.1, 2.25, 5.5, 0.25, "MODULE OUTCOME", 8, bold=True, color=LIGHT)
    _textbox(s, M + 0.1, 2.55, 5.5, 0.4, "Evaluate ethical frameworks", BODY_PT, bold=True, color=INK)
    box2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(2.2), Inches(5.8), Inches(0.9))
    box2.fill.solid()
    box2.fill.fore_color.rgb = BG
    box2.line.color.rgb = LIGHT
    _textbox(s, 6.6, 2.25, 5.5, 0.25, "ASSESSMENT", 8, bold=True, color=LIGHT)
    _textbox(s, 6.6, 2.55, 5.5, 0.4, "Multiple-choice quiz on definitions", BODY_PT, bold=True, color=INK)
    _textbox(s, M, 3.4, 11, 0.35, "What is the misalignment? Write your answer before the next slide.", SMALL_PT, italic=True, color=MID)
    _footer(s, 17)

    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 3, 0.25, "Phase 2 \u2014 Worked example reveal", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 8, 0.5, "The diagnosis", TITLE_PT, bold=True, color=INK)
    tbl2 = s.shapes.add_table(4, 2, Inches(M), Inches(1.1), Inches(12), Inches(2.8)).table
    tbl2.cell(0, 0).text = "OUTCOME DEMANDS"
    tbl2.cell(0, 1).text = "Evaluation \u2014 apply a schema to a novel ethical situation"
    tbl2.cell(1, 0).text = "ASSESSMENT TESTS"
    tbl2.cell(1, 1).text = "Recognition memory \u2014 recall of labels and definitions"
    tbl2.cell(2, 0).text = "GAP"
    tbl2.cell(2, 1).text = "Schema application requires accessible knowledge AND practiced judgment. The assessment only tests the first condition."
    tbl2.cell(3, 0).text = "COMMITMENT VIOLATED"
    tbl2.cell(3, 1).text = "Memory precedes complex thinking: the assessment never asks learners to think."
    _footer(s, 18)

    # --- Phase 3 ---
    _phase_header_slide(prs, "3", "Role Mapping",
        "Each role mediates learning differently. Making those differences concrete is the first step toward genuine coordination.", 19, image_placeholder=True)

    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 3, 0.25, "Phase 3 \u2014 Role clarification", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 9, 0.5, "Work within your role group", TITLE_PT, bold=True, color=INK)
    _textbox(s, M, 0.95, 11, 0.3, "Use the prompts in your workbook. 15 minutes.", SMALL_PT, color=MID)
    roles = [("Curriculum Design", "Alignment, sequencing, transfer"), ("Learning / Experience Design", "Cognitive operations, retrieval, visible reasoning"),
             ("Multimedia Design", "Cognitive function, load, schema support"), ("Learning Technology", "Affordances, constraints, analytics")]
    for i, (role, focus) in enumerate(roles):
        cx, cy = (M, 1.5 + (i // 2) * 1.4) if i % 2 == 0 else (6.8, 1.5 + (i // 2) * 1.4)
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(cy), Inches(5.8), Inches(0.75))
        rect.fill.solid()
        rect.fill.fore_color.rgb = BG
        rect.line.fill.background()
        _textbox(s, cx + 0.08, cy + 0.02, 5.5, 0.35, role, SMALL_PT, bold=True, color=INK)
        _textbox(s, cx + 0.08, cy + 0.32, 5.5, 0.35, focus, 8, color=MID)
    _footer(s, 20)

    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 4, 0.25, "Phase 3 \u2014 Cross-role perspective taking", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 9, 0.5, "Pair with a different role", TITLE_PT, bold=True, color=INK)
    _textbox(s, M, 1.1, 11, 0.4, "Step 1   Describe a real design decision you made recently.", BODY_PT, color=INK)
    _textbox(s, M, 1.65, 11, 0.5, "Step 2   Your partner identifies which learning science commitment it serves \u2014 or violates.", BODY_PT, color=MID)
    _textbox(s, M, 2.35, 11, 0.4, "Step 3   Switch. 2 minutes per person.", BODY_PT, color=MID)
    _footer(s, 21)

    s = prs.slides.add_slide(blank)
    trans_bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    trans_bg.fill.solid()
    trans_bg.fill.fore_color.rgb = DARK_BG
    trans_bg.line.fill.background()
    _textbox(s, M, 0.35, 3, 0.25, "TRANSITION", LABEL_PT, bold=True, color=RGBColor(136, 136, 136))
    _textbox(s, M, 1.5, 8, 0.5, "The question shifts.", 22, bold=True, color=WHITE)
    _textbox(s, M, 2.3, 11, 0.4, "From: \u201cWhat does my role do?\u201d", 16, color=RGBColor(170, 170, 170))
    _textbox(s, M, 2.8, 11, 0.4, "To: \u201cWhere is the system failing the learner?\u201d", 16, bold=True, color=WHITE)
    _textbox(s, M, 3.5, 11, 0.3, "Mixed-role team assignments will now be revealed.", SMALL_PT, color=RGBColor(102, 102, 102))
    nb2 = s.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.6), Inches(0.3))
    nb2.text_frame.paragraphs[0].text = "22 / 37"
    nb2.text_frame.paragraphs[0].font.color.rgb = RGBColor(136, 136, 136)

    # --- Phase 4 ---
    _phase_header_slide(prs, "4", "Cross-Role Diagnosis",
        "Apply learning science as a shared diagnostic lens to real design artefacts you brought to this room.", 23, image_placeholder=True)

    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 3, 0.25, "Phase 4 \u2014 Ground rule", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 8, 0.5, "Reminder before we begin", TITLE_PT, bold=True, color=INK)
    _textbox(s, M, 1.4, 11, 0.4, "\u201cWe\u2019re diagnosing the system, not individuals.", 16, italic=True, color=INK)
    _textbox(s, M, 1.85, 11, 0.4, "Misalignment is structural, not personal.", 16, italic=True, color=INK)
    _textbox(s, M, 2.3, 11, 0.4, "Our goal is coherence, not blame.\u201d", 16, italic=True, color=INK)
    _textbox(s, M, 3.0, 11, 0.35, "Point to specific elements in the artefact. Not general impressions.", SMALL_PT, color=LIGHT)
    _footer(s, 24)

    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 3, 0.25, "Phase 4 \u2014 Stage 1", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 8, 0.5, "Identify", TITLE_PT, bold=True, color=INK)
    _textbox(s, M, 0.95, 11, 0.3, "5 minutes. Point to specific elements in the artefact.", SMALL_PT, color=MID)
    for i, q in enumerate(["What must learners remember?", "Where is retrieval structured?", "Where is thinking visible?", "Where is cognitive load unnecessary?"], 1):
        _textbox(s, M, 1.5 + i * 0.5, 11, 0.45, f"{i}.  {q}", BODY_PT, color=INK)
    _footer(s, 25)

    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 3, 0.25, "Phase 4 \u2014 Stage 2", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 8, 0.5, "Diagnose misalignment", TITLE_PT, bold=True, color=INK)
    _textbox(s, M, 0.95, 11, 0.3, "20 minutes. Work through both questions.", SMALL_PT, color=MID)
    _textbox(s, M, 1.5, 11, 0.5, "5.  Where do role decisions contradict each other?", BODY_PT, bold=True, color=INK)
    _textbox(s, M, 1.85, 11, 0.35, "List at least two contradictions.", SMALL_PT, italic=True, color=MID)
    _textbox(s, M, 2.45, 11, 0.5, "6.  Which contradictions most impair learning?", BODY_PT, bold=True, color=INK)
    _textbox(s, M, 2.8, 11, 0.35, "Rank by impact on the learner. This is your redesign focus.", SMALL_PT, italic=True, color=MID)
    _footer(s, 26)

    s = prs.slides.add_slide(blank)
    break_bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    break_bg.fill.solid()
    break_bg.fill.fore_color.rgb = BG
    break_bg.line.fill.background()
    _textbox(s, 5.5, 3.2, 2.5, 0.6, "Break", 44, bold=True, color=RGBColor(204, 204, 204))
    _textbox(s, 5.8, 3.7, 2, 0.3, "10 minutes", BODY_PT, color=LIGHT)
    _textbox(s, 2.5, 4.2, 8.5, 0.4, "When you return: identify the one misalignment most blocking learning. That is your redesign focus.", SMALL_PT, italic=True, color=MID)
    s.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.6), Inches(0.3)).text_frame.paragraphs[0].text = "27 / 37"

    # --- Phase 5 ---
    _phase_header_slide(prs, "5", "Collaborative Redesign",
        "Translate the diagnosis into a coherent, learning-science-grounded redesign. Seven sections. One capability object. All roles contributing.", 28, image_placeholder=True)

    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 4, 0.25, "Phase 5 \u2014 Redesign map", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 9, 0.5, "Seven sections. One coherent object.", TITLE_PT, bold=True, color=INK)
    _textbox(s, M, 0.95, 11, 0.35, "Use the template in your workbook. 25 minutes. Coherent > complete.", SMALL_PT, color=MID)
    sections = [("1", "Capability Object"), ("2", "Retrieval Points"), ("3", "Visible Reasoning"), ("4", "Load Reduction"),
                 ("5", "Media Justification"), ("6", "Platform Alignment"), ("7", "Connection to Programme")]
    for i, (num, title) in enumerate(sections):
        cx = M + (i % 4) * 3.05
        cy = 1.5 + (i // 4) * 1.1
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(cy), Inches(2.9), Inches(0.95))
        rect.fill.solid()
        rect.fill.fore_color.rgb = BG
        rect.line.fill.background()
        _textbox(s, cx + 0.08, cy + 0.02, 2.7, 0.25, num + ".", LABEL_PT, bold=True, color=LIGHT)
        _textbox(s, cx + 0.08, cy + 0.28, 2.7, 0.6, title, SMALL_PT, bold=True, color=INK)
    _footer(s, 29)

    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 4, 0.25, "Phase 5 \u2014 Learning Science Justification", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 9, 0.5, "Before the gallery walk", TITLE_PT, bold=True, color=INK)
    _textbox(s, M, 0.95, 11, 0.4, "Complete the Learning Science Justification section. If you cannot answer these, the redesign is not yet grounded in science.", SMALL_PT, color=MID)
    for i, q in enumerate(["What memory structures are strengthened by this redesign?", "Where is retrieval embedded, and at what intervals?", "What cognitive load decisions were made? What was added, what was removed, and why?"], 1):
        _textbox(s, M, 1.55 + i * 0.55, 11, 0.5, q, BODY_PT, color=INK)
    _footer(s, 30)

    # --- Phase 6 ---
    _phase_header_slide(prs, "6", "Cross-Team Critique",
        "Structured feedback grounded in learning science. Three notes per redesign. No general impressions.", 31, image_placeholder=False)

    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 4, 0.25, "Phase 6 \u2014 Gallery walk protocol", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 9, 0.5, "Three notes. Every reviewer. Every redesign.", TITLE_PT, bold=True, color=INK)
    _textbox(s, M, 0.95, 11, 0.3, "Review at least two redesigns. 15 minutes.", SMALL_PT, color=MID)
    _textbox(s, M, 1.45, 11, 0.4, "Green sticky \u2014 Strength: What aligns well with learning science? Name the commitment it serves.", BODY_PT, bold=True, color=INK)
    _textbox(s, M, 1.85, 11, 0.4, "Yellow sticky \u2014 Tension: Where might the design impair learning? Name the commitment it violates.", BODY_PT, bold=True, color=INK)
    _textbox(s, M, 2.25, 11, 0.4, "Blue sticky \u2014 Clarification question: What needs more explanation to evaluate the design?", BODY_PT, bold=True, color=INK)
    _footer(s, 32)

    # --- Phase 7 ---
    _phase_header_slide(prs, "7", "Collaboration Charter",
        "Convert individual insight into concrete, accountable collective commitment.", 33, image_placeholder=False)

    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 4, 0.25, "Phase 7 \u2014 Individual commitments", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 9, 0.5, "Write before you share", TITLE_PT, bold=True, color=INK)
    _textbox(s, M, 0.95, 11, 0.35, "3 minutes. Individual writing first. No discussion until everyone has written.", SMALL_PT, color=MID)
    prompts = ["One decision I will no longer make alone", "One role I need earlier in the design process", "One learning science principle I will use in future discussions", "One process change we will implement in the next design cycle"]
    for i, p in enumerate(prompts):
        _textbox(s, M, 1.6 + i * 0.45, 11, 0.4, "\u2014  " + p, BODY_PT, color=INK)
    _footer(s, 34)

    s = prs.slides.add_slide(blank)
    _textbox(s, M, 0.35, 4, 0.25, "Phase 7 \u2014 30-day follow-up", LABEL_PT, bold=True, color=LIGHT)
    _textbox(s, M, 0.55, 9, 0.5, "In 4 weeks.", TITLE_PT, bold=True, color=INK)
    _textbox(s, M, 0.95, 11, 0.35, "We meet for 30 minutes. Each role shares three things.", BODY_PT, color=MID)
    _textbox(s, M, 1.5, 11, 0.4, "1.  One decision you made differently since the workshop", BODY_PT, color=INK)
    _textbox(s, M, 1.95, 11, 0.4, "2.  One place where collaboration improved", BODY_PT, color=INK)
    _textbox(s, M, 2.4, 11, 0.4, "3.  One remaining misalignment that needs addressing", BODY_PT, color=INK)
    _textbox(s, M, 3.2, 11, 0.35, "Calendar invitations go out today. This is not optional.", SMALL_PT, italic=True, color=LIGHT)
    _footer(s, 35)

    s = prs.slides.add_slide(blank)
    close_bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    close_bg.fill.solid()
    close_bg.fill.fore_color.rgb = DARK_BG
    close_bg.line.fill.background()
    _textbox(s, M, 0.35, 3, 0.25, "CLOSING RETRIEVAL TASK", LABEL_PT, bold=True, color=RGBColor(136, 136, 136))
    _textbox(s, M, 0.7, 8, 0.5, "No notes. Write three things.", 28, bold=True, color=WHITE)
    _textbox(s, M, 1.6, 11, 0.4, "\u2014  How you now define learning", 15, color=RGBColor(204, 204, 204))
    _textbox(s, M, 2.05, 11, 0.4, "\u2014  One way your role mediates cognitive change", 15, color=RGBColor(204, 204, 204))
    _textbox(s, M, 2.5, 11, 0.4, "\u2014  One collaboration commitment you are taking forward", 15, color=RGBColor(204, 204, 204))
    _textbox(s, M, 3.2, 11, 0.3, "2 minutes. Silent. These are for you, not for the group.", SMALL_PT, color=RGBColor(102, 102, 102))
    nb3 = s.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.6), Inches(0.3))
    nb3.text_frame.paragraphs[0].text = "36 / 37"
    nb3.text_frame.paragraphs[0].font.color.rgb = RGBColor(136, 136, 136)

    # --- Slide 37: Close ---
    s = prs.slides.add_slide(blank)
    left_white = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(5.6), Inches(7.5))
    left_white.fill.solid()
    left_white.fill.fore_color.rgb = WHITE
    left_white.line.fill.background()
    right_dark = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.6), 0, Inches(7.733), Inches(7.5))
    right_dark.fill.solid()
    right_dark.fill.fore_color.rgb = DARK_BG
    right_dark.line.fill.background()
    _placeholder_rect(s, 0, 0, 5.4, 7.5, "[ CLOSING IMAGE ]")
    _textbox(s, 5.8, 0.35, 7, 0.25, "YOU LEAVE WITH", LABEL_PT, bold=True, color=RGBColor(102, 102, 102))
    for i, t in enumerate(["A redesigned learning experience artefact", "A shared operational definition of learning", "Practical learning science commitments", "A cross-role collaboration charter", "A 30-day follow-up plan"]):
        _textbox(s, 5.8, 0.75 + i * 0.35, 7, 0.35, "\u2014  " + t, SMALL_PT, color=WHITE)
    nb4 = s.shapes.add_textbox(Inches(12.5), Inches(7.0), Inches(0.6), Inches(0.3))
    nb4.text_frame.paragraphs[0].text = f"37 / {TOTAL_SLIDES}"
    nb4.text_frame.paragraphs[0].font.color.rgb = RGBColor(136, 136, 136)

    prs.save(output_path)
    print(f"Slides written to: {output_path}")


if __name__ == "__main__":
    out = os.path.join(os.path.dirname(__file__), "workshop-slides.pptx")
    build_pptx(out)
